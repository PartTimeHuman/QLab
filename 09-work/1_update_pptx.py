#!/usr/bin/env python3
"""
1_update_pptx.py
烧碱周报自动化更新脚本
==================================================
功能:
1. 将CSV数据库转换为结构化Excel（多Sheet，按分类整理）
2. 自动更新PPT中的图表数据（季节性曲线图）
3. 自动更新PPT中的文字注释（使用{占位符}语法）

使用方法:
    python 1_update_pptx.py

依赖:
    pip install pandas openpyxl python-pptx lxml
"""

import sys
import warnings
from pathlib import Path
from copy import deepcopy
from datetime import datetime
import re

import pandas as pd
import numpy as np

warnings.filterwarnings("ignore")

# ─────────────────────────────────────────────────────────
# 配置区（用户可按需修改）
# ─────────────────────────────────────────────────────────

# 脚本所在目录（文件放在同一个文件夹即可自动找到）
HERE = Path(__file__).parent.resolve()

# 输入文件（与本脚本放在同一目录）
CSV_FILE      = HERE / "工作簿2.csv"
TEMPLATE_PPTX = HERE / "烧碱周报模版.pptx"

# 输出文件（也输出到同一目录）
OUTPUT_EXCEL  = HERE / "烧碱数据库.xlsx"
OUTPUT_PPTX   = HERE / f"烧碱周报_{datetime.today().strftime('%Y%m%d')}.pptx"

# ─────────────────────────────────────────────────────────
# 图表配置：PPT中每个图表对应哪个数据列
# 说明：key = (幻灯片编号, 图表shape名称)，value = 数据库列名
# 你可以通过运行"inspect_charts()"函数来发现PPT中所有图表的名称
# ─────────────────────────────────────────────────────────
CHART_COLUMN_MAP = {
    # 幻灯片8 - 现货
    (8, "图表 21"): "山东32交割库最低价",
    (8, "图表 22"): "江苏32交割库最低价",
    (8, "图表 24"): "山东50交割库最低价",
    (8, "图表 29"): "魏桥采购价",
    (8, "图表 30"): "华南50",
    (8, "图表 31"): "液氯价格-山东",
    # 幻灯片9 - 基差
    (9, "图表 14"): "山东32最低交割基差",
    (9, "图表 15"): "江苏32最低交割基差",
    (9, "图表 25"): "浙江32最低交割基差",
    (9, "图表 26"): "魏桥32基差",
    # 幻灯片10 - 月差
    (10, "图表 10"): "1-5月差",
    (10, "图表 11"): "5-9月差",
    (10, "图表 12"): "9-1月差",
    # 幻灯片13 - 产量
    (13, "图表 15"): "烧碱周产量(wt)",
    (13, "图表 16"): "华东烧碱周产(wt)",
    (13, "图表 20"): "华北烧碱周产(wt)",
    # 幻灯片14 - 山东库存
    (14, "图表 11"): "山东32库存-湿吨",
    (14, "图表 12"): "山东50库存湿吨",
    # 幻灯片16 - 成本利润
    (16, "图表 13"): "山东氯碱利润-完全成本",
    (16, "图表 21"): "盘面边际利润",
    # 幻灯片22 - 出口
    (22, "图表 22"): "烧碱总出口(wt)",
    (22, "图表 23"): "烧碱出口：印尼",
    (22, "图表 24"): "烧碱出口：澳大利亚",
    # 幻灯片37 - 平衡表
    (37, "图表 8"):  "产量-平衡表",
    (37, "图表 9"):  "总需求-平衡表",
    (37, "图表 10"): "库存-平衡表",
}

# ─────────────────────────────────────────────────────────
# 文字注释配置：PPT中需要自动更新的文本框
# PPT里的文字框文本内包含 {变量名} 时会被自动替换
# 格式: {col__func} 例如 {烧碱总出口(wt)__last} 取最新值
#       支持的函数: last(最新值) yoy(同比%) mom(环比%) ytd(年初至今累计)
# ─────────────────────────────────────────────────────────
# 在PPT文字框中直接写这样的标记:
# "全国烧碱本周出口{烧碱总出口(wt)__last}万吨，同比{烧碱总出口(wt)__yoy}%"
# 脚本会自动替换成最新数值


# ─────────────────────────────────────────────────────────
# 第一步：加载数据库
# ─────────────────────────────────────────────────────────
def load_database(filepath: str) -> pd.DataFrame:
    """加载CSV数据库，返回 DatetimeIndex × 指标名 的DataFrame"""
    print(f"[1/4] 加载数据库: {filepath}")
    df_raw = pd.read_csv(filepath, low_memory=False)
    real_cols = [str(c).strip() for c in df_raw.iloc[0].tolist()]
    real_cols[0] = "日期"

    # 去重
    seen = {}
    unique_cols = []
    for c in real_cols:
        if c in seen:
            seen[c] += 1
            unique_cols.append(f"{c}_{seen[c]}")
        else:
            seen[c] = 0
            unique_cols.append(c)

    df_raw.columns = unique_cols
    df = df_raw.iloc[1:].copy()
    df = df.set_index("日期")
    df.index = pd.to_datetime(df.index, errors="coerce")
    df = df[df.index.notna()].sort_index()
    df = df.apply(pd.to_numeric, errors="coerce")
    df = df.dropna(axis=1, how="all")

    print(f"      → {df.shape[0]} 行 × {df.shape[1]} 列  "
          f"({df.index[0].date()} ~ {df.index[-1].date()})")
    return df


# ─────────────────────────────────────────────────────────
# 第二步：导出结构化Excel
# ─────────────────────────────────────────────────────────
EXCEL_SHEETS = {
    "现货估值": [
        "山东32交割库最低价", "山东50交割库最低价", "江苏32交割库最低价",
        "魏桥采购价", "东营32价格", "西北片碱均价", "华南50", "出口FOB",
        "液氯价格-山东",
        "山东32最低交割基差", "山东32最低出厂基差", "江苏32最低交割基差",
        "浙江32最低交割基差", "魏桥32基差", "山东50最低交割基差",
        "1-5月差", "5-9月差", "9-1月差",
        "50-32价差：华泰", "50-32价差：金岭",
        "日波动率", "年化波动率",
    ],
    "供给": [
        "烧碱周产量(wt)", "烧碱月产量(wt)", "华北烧碱月产量(wt)",
        "华北烧碱周产(wt)", "华东烧碱周产(wt)", "西北烧碱周产(wt)",
        "山东32碱产量", "山东50产量", "山东日产量",
        "月度开工率", "周度开工率", "月度装置损失率", "周度装置损失率",
        "山东液碱开工率",
        "烧碱总库存(wt)", "液碱总库存(wt)", "固碱总库存(wt)",
        "山东32库存-湿吨", "山东50库存湿吨", "山东总库存-湿吨",
        "魏桥库存-湿吨", "烧碱仓单",
        "山东氯碱利润-完全成本", "山东氯碱利润-现金流成本",
        "盘面边际利润", "01盘面利润", "05盘面利润", "09盘面利润",
        "魏桥价格边际利润", "山东原盐价格", "液氯价格-山东",
    ],
    "需求": [
        "浙江印染开机", "江苏印染开机", "华东印染开机",
        "粘胶短纤产量(周)(万吨)", "粘胶短纤产能利用率(周)",
        "粘胶短纤产量(月)(万吨)", "粘胶短纤产能利用率(月)",
        "铝土矿产量", "铝土矿港口库存", "氧化铝产量", "氧化铝产能利用率",
        "氧化铝利润", "氧化铝净进口", "氧化铝总库存",
        "聚合MDI产量(万吨)", "聚合MDI产能利用率",
        "纯MDI产量(万吨)", "纯MDI产能利用率",
        "PC产量(wt)", "环氧树脂产量(wt)",
        "烧碱总出口(wt)", "烧碱出口单价(USD/t)", "净出口(wt)",
        "烧碱出口：澳大利亚", "烧碱出口：印尼", "烧碱出口：台湾",
        "烧碱出口：越南", "烧碱出口：日本", "烧碱出口：韩国", "烧碱出口：印度",
    ],
    "平衡表": [
        "产量-平衡表", "总需求-平衡表", "库存-平衡表", "出口-平衡表",
        "总需求-平衡表（实产法）", "库存-平衡表（实产法）",
    ],
    "联产利润": [
        "环氧丙烷利润-外采氯", "环氧丙烷利润-自用氯", "环氧丙烷完全利润",
        "环氧丙烷：产能利用率：中国（周）",
        "ECH利润-甘油法", "ECH利润-外采氯丙烯法",
        "甲烷氯化物利润-外采氯", "甲烷氯化物利润-自用氯",
        "PVC现货利润含烧碱", "V01盘面利润含烧碱", "V05盘面利润含烧碱",
        "烧碱-环氧丙烷联产利润", "烧碱-ECH联产利润-甘油法",
        "烧碱-ECH联产利润-丙烯法", "烧碱-甲烷氯化物联产利润",
        "R22：生产毛利", "软泡聚醚利润", "硬泡聚醚利润",
    ],
}


def export_excel(df: pd.DataFrame, output_path: str):
    """将DataFrame按分类导出到多Sheet Excel文件"""
    print(f"[2/4] 导出结构化Excel: {output_path}")
    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        # Sheet 1: 所有数据
        df.to_excel(writer, sheet_name="ALL_DATA", freeze_panes=(1, 1))

        # Sheet 2: 元数据（各列数据量、最新值）
        meta = []
        for col in df.columns:
            s = df[col].dropna()
            if len(s) == 0:
                continue
            meta.append({
                "指标名称": col,
                "有效数据量": len(s),
                "起始日期": s.index[0].date(),
                "最新日期": s.index[-1].date(),
                "最新值": round(s.iloc[-1], 4),
                "均值": round(s.mean(), 4),
                "最大值": round(s.max(), 4),
                "最小值": round(s.min(), 4),
            })
        pd.DataFrame(meta).to_excel(
            writer, sheet_name="元数据", index=False, freeze_panes=(1, 0)
        )

        # 分类Sheet
        for sheet_name, cols in EXCEL_SHEETS.items():
            available = [c for c in cols if c in df.columns]
            if not available:
                continue
            df[available].to_excel(
                writer, sheet_name=sheet_name, freeze_panes=(1, 1)
            )
            print(f"      → Sheet [{sheet_name}]: {len(available)} 列")

    print(f"      ✓ Excel导出完成")


# ─────────────────────────────────────────────────────────
# 第三步：更新PPT图表
# ─────────────────────────────────────────────────────────
def build_seasonal_series(series: pd.Series, years: list) -> dict:
    """
    将时间序列拆分为按年的季节性数据。
    返回: {year: [366个值（含NaN填充）]} 
    day_index 从 0 开始（对应1月1日）
    """
    result = {}
    for year in years:
        year_data = series[series.index.year == year]
        # 创建全年366个槽位（兼容闰年）
        arr = [np.nan] * 366
        for ts, val in year_data.items():
            doy = ts.dayofyear - 1  # 0-indexed
            if doy < 366:
                arr[doy] = None if pd.isna(val) else float(val)
        result[year] = arr
    return result


def update_chart_in_pptx(slide, shape_name: str, col_name: str, df: pd.DataFrame):
    """更新PPT中指定图表的数据（季节性年度对比图）"""
    from pptx.chart.data import ChartData
    from lxml import etree

    # 找到目标shape
    target = None
    for shape in slide.shapes:
        if shape.name == shape_name and shape.shape_type == 3:
            target = shape
            break

    if target is None:
        print(f"      ✗ 未找到图表: {shape_name}")
        return False

    if col_name not in df.columns:
        print(f"      ✗ 数据列不存在: {col_name}")
        return False

    series = df[col_name].dropna()
    if len(series) == 0:
        return False

    chart = target.chart

    # 确定要更新的年份（最近6年 + 当前年）
    current_year = datetime.today().year
    years = list(range(current_year - 5, current_year + 1))

    # 尝试使用 python-pptx 的 replace_data 方法
    try:
        seasonal = build_seasonal_series(series, years)

        chart_data = ChartData()
        # 用1-366作为横坐标
        chart_data.categories = list(range(1, 367))

        for year in years:
            vals = seasonal.get(year, [None] * 366)
            chart_data.add_series(str(year), vals)

        chart.replace_data(chart_data)
        print(f"      ✓ 更新图表 [{shape_name}] ← {col_name}")
        return True

    except Exception as e:
        print(f"      ✗ 更新图表 [{shape_name}] 失败: {e}")
        return False


# ─────────────────────────────────────────────────────────
# 第四步：更新PPT文字注释
# ─────────────────────────────────────────────────────────
def compute_placeholder_value(col_name: str, func: str, df: pd.DataFrame) -> str:
    """计算占位符的值"""
    if col_name not in df.columns:
        return f"[{col_name}:无数据]"

    s = df[col_name].dropna()
    if len(s) == 0:
        return "[无数据]"

    today = datetime.today()

    if func == "last":
        return f"{s.iloc[-1]:.2f}"

    elif func == "yoy":
        last_val = s.iloc[-1]
        last_date = s.index[-1]
        # 找去年同期
        one_year_ago = last_date - pd.DateOffset(years=1)
        nearby = s[abs(s.index - one_year_ago) <= pd.Timedelta("7d")]
        if len(nearby) == 0:
            return "N/A"
        yoy_val = nearby.iloc[-1]
        pct = (last_val - yoy_val) / abs(yoy_val) * 100
        sign = "+" if pct >= 0 else ""
        return f"{sign}{pct:.1f}%"

    elif func == "mom":
        if len(s) < 2:
            return "N/A"
        last_val = s.iloc[-1]
        prev_val = s.iloc[-2]
        pct = (last_val - prev_val) / abs(prev_val) * 100
        sign = "+" if pct >= 0 else ""
        return f"{sign}{pct:.1f}%"

    elif func == "ytd":
        year_start = pd.Timestamp(today.year, 1, 1)
        ytd_data = s[s.index >= year_start]
        if len(ytd_data) == 0:
            return "N/A"
        return f"{ytd_data.sum():.2f}"

    elif func == "ytd_yoy":
        year_start = pd.Timestamp(today.year, 1, 1)
        last_date = s.index[-1]
        ytd_curr = s[(s.index >= year_start) & (s.index <= last_date)].sum()
        ytd_prev = s[
            (s.index >= year_start - pd.DateOffset(years=1)) &
            (s.index <= last_date - pd.DateOffset(years=1))
        ].sum()
        if ytd_prev == 0:
            return "N/A"
        pct = (ytd_curr - ytd_prev) / abs(ytd_prev) * 100
        sign = "+" if pct >= 0 else ""
        return f"{sign}{pct:.1f}%"

    return f"[未知函数:{func}]"


PLACEHOLDER_PATTERN = re.compile(r"\{([^}]+)__([^}]+)\}")


def update_text_in_pptx(prs, df: pd.DataFrame):
    """
    扫描PPT所有文字框，将 {列名__函数} 形式的占位符替换为实际数值。
    支持的函数: last, yoy, mom, ytd, ytd_yoy
    """
    replaced_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    if "{" not in text:
                        continue

                    def replace_match(m):
                        col_name, func = m.group(1), m.group(2)
                        return compute_placeholder_value(col_name, func, df)

                    new_text = PLACEHOLDER_PATTERN.sub(replace_match, text)
                    if new_text != text:
                        run.text = new_text
                        replaced_count += 1

    print(f"      ✓ 文字占位符更新: {replaced_count} 处")
    return replaced_count


# ─────────────────────────────────────────────────────────
# 辅助工具：列出PPT所有图表
# ─────────────────────────────────────────────────────────
def inspect_charts(pptx_path: str):
    """打印PPT中所有图表的位置和名称，方便填写CHART_COLUMN_MAP"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    print("\n=== PPT图表清单 ===")
    print(f"{'幻灯片':>6} | {'图表名称':<20} | {'图表类型'}")
    print("-" * 50)
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 3:
                try:
                    chart_type = shape.chart.chart_type
                except:
                    chart_type = "未知"
                print(f"  {slide_idx:>4} | {shape.name:<20} | {chart_type}")


# ─────────────────────────────────────────────────────────
# 主流程
# ─────────────────────────────────────────────────────────
def main():
    from pptx import Presentation

    print("=" * 60)
    print("烧碱周报自动化更新脚本")
    print("=" * 60)

    # 检查文件
    for f in [CSV_FILE, TEMPLATE_PPTX]:
        if not Path(f).exists():
            print(f"错误: 找不到文件 {f}")
            sys.exit(1)

    # Step 1: 加载数据库
    df = load_database(CSV_FILE)

    # Step 2: 导出Excel
    export_excel(df, OUTPUT_EXCEL)

    # Step 3: 更新PPT
    print(f"[3/4] 更新PPT图表: {TEMPLATE_PPTX}")
    prs = Presentation(TEMPLATE_PPTX)

    slides_list = list(prs.slides)
    chart_update_count = 0

    for (slide_no, chart_name), col_name in CHART_COLUMN_MAP.items():
        if slide_no < 1 or slide_no > len(slides_list):
            print(f"      ✗ 幻灯片编号越界: {slide_no}")
            continue
        slide = slides_list[slide_no - 1]
        ok = update_chart_in_pptx(slide, chart_name, col_name, df)
        if ok:
            chart_update_count += 1

    print(f"      → 图表更新完成: {chart_update_count}/{len(CHART_COLUMN_MAP)} 个")

    # Step 4: 更新文字注释
    print(f"[4/4] 更新PPT文字注释")
    update_text_in_pptx(prs, df)

    # 保存
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ 完成！输出文件:")
    print(f"   Excel: {OUTPUT_EXCEL}")
    print(f"   PPT:   {OUTPUT_PPTX}")
    print()
    print("💡 提示: 若要添加更多图表映射，编辑 CHART_COLUMN_MAP 字典")
    print("💡 提示: 在PPT文字框中使用 {列名__last} 等占位符实现自动文字更新")
    print("💡 提示: 运行 inspect_charts('烧碱周报模版.pptx') 查看所有图表名称")


if __name__ == "__main__":
    # 如需查看PPT图表清单，取消下行注释:
    # inspect_charts(TEMPLATE_PPTX)

    main()
